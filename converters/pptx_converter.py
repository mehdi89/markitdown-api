from pptx import Presentation
import re
import logging
import os

logger = logging.getLogger(__name__)

def convert_pptx_to_markdown(file_path, use_llm=False):
    """
    Convert PowerPoint file to Markdown
    
    Args:
        file_path (str): Path to the PowerPoint file
        use_llm (bool): Whether to use LLM enhancement (not implemented)
        
    Returns:
        str: Markdown content
    """
    logger.debug(f"Converting PowerPoint file: {file_path}")
    
    try:
        # Get the filename without extension for use as title
        file_name = os.path.basename(file_path)
        file_name_without_ext = os.path.splitext(file_name)[0]
        
        # Open the presentation
        prs = Presentation(file_path)
        
        # Start with a title
        markdown_text = f"# {file_name_without_ext}\n\n"
        
        # Process each slide
        for slide_number, slide in enumerate(prs.slides, 1):
            markdown_text += f"## Slide {slide_number}\n\n"
            
            # Process slide title
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                markdown_text += f"### {title_text}\n\n"
            
            # Process text elements in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                    text = shape.text.strip()
                    
                    # Try to detect if it's a bullet list
                    lines = text.split('\n')
                    if any(line.strip().startswith('•') for line in lines):
                        for line in lines:
                            line = line.strip()
                            if line.startswith('•'):
                                markdown_text += f"* {line[1:].strip()}\n"
                            else:
                                markdown_text += f"{line}\n"
                        markdown_text += "\n"
                    else:
                        markdown_text += f"{text}\n\n"
            
            # Add separator between slides
            markdown_text += "---\n\n"
        
        # Clean up multiple newlines
        markdown_text = re.sub(r'\n{3,}', '\n\n', markdown_text)
        
        if use_llm:
            logger.info("LLM enhancement requested but not implemented for PowerPoint")
        
        return markdown_text
    
    except Exception as e:
        logger.error(f"Error converting PowerPoint to Markdown: {str(e)}")
        raise Exception(f"PowerPoint conversion error: {str(e)}")